from pptx import Presentation
from pptx.util import Pt
import copy
import six
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from bs4 import BeautifulSoup
import os
import sys
import argparse
from datetime import datetime

updatePPTX=1.10

def parse_percentage(text):
    try:
        # Récupère uniquement le premier nombre dans la chaîne
        import re
        match = re.search(r'^(\d+\.?\d*)', text.strip())
        if match:
            # Convertit en float le premier groupe capturé
            return float(match.group(1))
        return None
    except (ValueError, AttributeError):
        # En cas d'erreur, retourne None
        return None
    
def add_new_slide(pres, date, clusterviewmode=False):
    # on recopie la derniere slide présente car c'est toujours la up to date slide
    source_slide = pres.slides[-1]

    # Blank slide layout is usually the last layout in the list
    try:
        blank_slide_layout = pres.slide_layouts[-1]
    except:
        blank_slide_layout = pres.slide_layouts[len(pres.slide_layouts)]

    new_slide = pres.slides.add_slide(blank_slide_layout)

    import io

    for shape in source_slide.shapes:
        el = shape.element
        if not "image" in dir(shape): 
            newel = copy.deepcopy(el)
            new_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')
        if "image" in dir(shape):
            img = io.BytesIO(shape.image.blob)
            new_shape_img = new_slide.shapes.add_picture(image_file = img,
                                            left = shape.left,
                                            top = shape.top,
                                            width = shape.width,
                                            height = shape.height)
            new_shape_img.name = shape.name
            #print(f"Image copied: {shape.name}")
            
    for shape in new_slide.shapes:
        if shape.name == "date_slide":
            if shape.text_frame.paragraphs[0].font.size:
                font_size = shape.text_frame.paragraphs[0].font.size
            else:
                font_size = Pt(12)
            if clusterviewmode:
                shape.text_frame.text = str(date)+"\nclusterview"
            else:  
                shape.text_frame.text = str(date)
            shape.text_frame.paragraphs[0].font.size = font_size
            shape.text_frame.paragraphs[0].font.bold = True
            shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
    
    # We need to copy the relationships from the source slide to the new slide

    for _, value in six.iteritems(source_slide.part.rels):
        # Make sure we don't copy a notesSlide relation as that won't exist
        if "notesSlide" not in value.reltype:
            new_slide.part.rels.get_or_add(
                value.reltype,
                value._target
            )

    return new_slide

def modif_textbox(shape,value,textvalue,warning,error):
    if shape.text_frame.paragraphs[0].font.size:
        font_size = shape.text_frame.paragraphs[0].font.size
    else:
        font_size = Pt(6)  
    shape.text_frame.text = str(textvalue+'%')
    shape.text_frame.paragraphs[0].font.size = font_size
    color_text = None
    if not value == None:
        if value >=warning:
            color_text= RGBColor(255, 255, 0)  # Yellow
        if value >=error:    
            color_text= RGBColor(192, 0, 0)  # Dark Red
        shape.text_frame.paragraphs[0].font.bold = True
        if color_text:
            shape.text_frame.paragraphs[0].font.color.rgb = color_text    

def update_shapes_in_slide(slide, table_data, replace_inplace=False):
    """
    Met à jour les shapes dans une slide donnée avec les données du tableau HTML
    """
    updates_count = 0
    
    # Parcourir les lignes du tableau
    for row in table_data.find_all('tr'):
        if row.next.name == 'th':
            continue        
        cells = row.find_all(['td', 'th'])
        
        # Vérifier qu'il y a 12 colonnes
        if len(cells) >= 13:
            nodename = cells[4].get_text().strip()
            capacity_percentage = cells[9].get_text().strip()
            headroom_percentage = cells[11].get_text().strip()
            float_value_capacity = parse_percentage(capacity_percentage)
            float_value_headroom = parse_percentage(headroom_percentage)
            
            # Rechercher la TextBox correspondante et mettre à jour le texte si necessaire
            for shape in slide.shapes:
                #print(f"Shape found: {shape.name}")
                if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    #print("Shape is a group, digging deeper...")
                    for sub_shape in shape.shapes:
                        #print(f"\tSub Shape found: {sub_shape.name}")
                        if sub_shape.name.lower() == str("capacity_" + nodename).lower() and sub_shape.has_text_frame:
                            sub_shape_value = parse_percentage(sub_shape.text_frame.text)
                            if sub_shape_value != float_value_capacity:
                                modif_textbox(sub_shape, float_value_capacity, str(float_value_capacity), 70, 80)
                                mode = "inplace" if replace_inplace else "new slide"
                                print(f"\tUpdate capacity ({mode}): {nodename} from {str(sub_shape_value)}% => {str(float_value_capacity)}%")
                                updates_count += 1
                            
                        if sub_shape.name.lower() == str("headroom_" + nodename).lower() and sub_shape.has_text_frame:
                            sub_shape_value = parse_percentage(sub_shape.text_frame.text)
                            if sub_shape_value != float_value_headroom:
                                modif_textbox(sub_shape, float_value_headroom, str(float_value_headroom), 80, 90)
                                mode = "inplace" if replace_inplace else "new slide"
                                print(f"\tUpdate headroom ({mode}): {nodename} from {str(sub_shape_value)}% => {str(float_value_headroom)}%")
                                updates_count += 1
                        
                if shape.name.lower() == str("capacity_" + nodename).lower() and shape.has_text_frame:
                    shape_value = parse_percentage(shape.text_frame.text)
                    if shape_value != float_value_capacity:
                        modif_textbox(shape, float_value_capacity, str(float_value_capacity), 70, 80)
                        mode = "inplace" if replace_inplace else "new slide"
                        print(f"Update capacity ({mode}): {nodename} from {str(shape_value)}% => {str(float_value_capacity)}%")
                        updates_count += 1
                    
                if shape.name.lower() == str("headroom_" + nodename).lower() and shape.has_text_frame:
                    shape_value = parse_percentage(shape.text_frame.text)
                    if shape_value != float_value_headroom:
                        modif_textbox(shape, float_value_headroom, str(float_value_headroom), 80, 90)
                        mode = "inplace" if replace_inplace else "new slide"
                        print(f"Update headroom ({mode}): {nodename} from {str(shape_value)}% => {str(float_value_headroom)}%")
                        updates_count += 1
                    
        else:
            print(f"row ignored, not enough columns: {row.get_text(strip=True)}")
    
    return updates_count

def update_ppt_from_html(ppt_path, html_path, replace_inplace=False, clusterviewmode=False):
    # Charger la présentation PowerPoint
    prs = Presentation(ppt_path)
    
    # Lire et parser le fichier HTML
    with open(html_path, 'r', encoding='utf-8') as f:
        html_content = f.read()
    
    # Recuperer la date de modification du fichier HTML
    file_time = os.path.getmtime(html_path)
    file_date = datetime.fromtimestamp(file_time).strftime('%m/%d/%Y')

    # Parser le contenu HTML pour trouver le tableau
    soup = BeautifulSoup(html_content, 'html.parser')
    table = soup.find('table')
    
    # Vérifier si un tableau a été trouvé
    if not table:
        raise ValueError("Aucun tableau trouvé dans le fichier HTML")
    
    total_updates = 0
    
    if replace_inplace:
        # Mode remplacement sur place : modifier toutes les slides existantes
        print(f"Mode replace_inplace activé - Modification de toutes les slides existantes...")
        
        for slide_idx, slide in enumerate(prs.slides):
            print(f"Traitement de la slide {slide_idx + 1}/{len(prs.slides)}...")
            updates_count = update_shapes_in_slide(slide, table, replace_inplace=True)
            total_updates += updates_count
            
        print(f"Total des mises à jour effectuées sur toutes les slides : {total_updates}")
        
    else:
        # Mode par défaut : ajouter une nouvelle slide
        print("Mode par défaut - Ajout d'une nouvelle slide...")
        new_slide = add_new_slide(prs, file_date, clusterviewmode=clusterviewmode)
        total_updates = update_shapes_in_slide(new_slide, table, replace_inplace=False)
        print(f"Nouvelle slide créée avec {total_updates} mises à jour")
    
    # Sauvegarder la présentation
    prs.save(ppt_path)
    mode_text = "modifié sur place" if replace_inplace else "avec nouvelle slide"
    print(f"\nPresentation updated and saved as: {ppt_path} ({mode_text})")

def main():
    # Configuration des arguments de ligne de commande
    parser = argparse.ArgumentParser(
        description="Met à jour une présentation PowerPoint avec les données d'un fichier HTML"
    )
    
    parser.add_argument(
        "ppt_file",
        help="Chemin vers le fichier PowerPoint à modifier"
    )
    
    parser.add_argument(
        "html_file", 
        help="Chemin vers le fichier HTML contenant les données"
    )
    
    parser.add_argument(
        "--replace-inplace",
        action="store_true",
        help="Modifier les slides existantes au lieu de créer une nouvelle slide"
    )
    
    parser.add_argument(
        "--clusterviewmode",
        action="store_true",
        help="Capacity via le mode Clusterview"
    )
    
    parser.add_argument(
        "--version",
        action="version",
        version=f"updatePPTX version '{updatePPTX}'"
    )

    # Parser les arguments
    args = parser.parse_args()
    
    # Vérifier que les fichiers existent
    if not os.path.exists(args.ppt_file):
        print(f"Erreur: Le fichier PowerPoint '{args.ppt_file}' n'existe pas.")
        sys.exit(1)
        
    if not os.path.exists(args.html_file):
        print(f"Erreur: Le fichier HTML '{args.html_file}' n'existe pas.")
        sys.exit(1)
    
    # Afficher les informations de version et de mode
    print(f"updatePPTX version '{updatePPTX}'")
    mode = "remplacement sur place" if args.replace_inplace else "ajout nouvelle slide"
    print(f"Mode: {mode}")
    print(f"PowerPoint: {args.ppt_file}")
    print(f"HTML: {args.html_file}")
    print()
    
    # Exécuter la mise à jour
    try:
        update_ppt_from_html(
            args.ppt_file, 
            args.html_file, 
            args.replace_inplace,
            args.clusterviewmode
        )
    except Exception as e:
        print(f"Erreur lors de la mise à jour: {e}")
        sys.exit(1)

if __name__ == "__main__":
    main()
